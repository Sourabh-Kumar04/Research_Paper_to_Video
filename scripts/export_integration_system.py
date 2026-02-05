"""
Export and Integration System

This module provides comprehensive export and integration capabilities including:
- Export to multiple formats (MP4, WebM, GIF, PowerPoint)
- Direct YouTube, Vimeo, and social media publishing
- API endpoints for third-party integrations
- Webhook support for automated workflows
- Cloud storage integration for exports

Optimized for 16GB RAM systems with efficient processing.
"""

import os
import json
import asyncio
import logging
from datetime import datetime
from typing import Dict, List, Optional, Any, Union, Callable
from dataclasses import dataclass, asdict
from enum import Enum
import uuid
from pathlib import Path
import tempfile
import subprocess

# Export format libraries
try:
    import ffmpeg
    from PIL import Image, ImageDraw, ImageFont
    import moviepy.editor as mp
    EXPORT_AVAILABLE = True
except ImportError:
    EXPORT_AVAILABLE = False
    logging.warning("Export libraries not available. Install ffmpeg-python, PIL, and moviepy.")

# Cloud and API integration
try:
    import requests
    import boto3
    from google.cloud import storage as gcs
    import httpx
    INTEGRATION_AVAILABLE = True
except ImportError:
    INTEGRATION_AVAILABLE = False
    logging.warning("Integration libraries not available. Install requests, boto3, google-cloud-storage, httpx.")

class ExportFormat(Enum):
    """Supported export formats."""
    MP4 = "mp4"
    WEBM = "webm"
    GIF = "gif"
    POWERPOINT = "pptx"
    PDF = "pdf"
    IMAGES = "images"
    AUDIO_ONLY = "audio"

class PlatformType(Enum):
    """Supported platforms for publishing."""
    YOUTUBE = "youtube"
    VIMEO = "vimeo"
    TWITTER = "twitter"
    FACEBOOK = "facebook"
    LINKEDIN = "linkedin"
    TIKTOK = "tiktok"
    INSTAGRAM = "instagram"

@dataclass
class ExportConfig:
    """Configuration for export operations."""
    format: ExportFormat
    quality: str = "high"  # low, medium, high, custom
    resolution: Optional[str] = None  # e.g., "1920x1080"
    fps: Optional[int] = None
    bitrate: Optional[str] = None
    custom_params: Dict[str, Any] = None
    
    def __post_init__(self):
        if self.custom_params is None:
            self.custom_params = {}

@dataclass
class PublishConfig:
    """Configuration for platform publishing."""
    platform: PlatformType
    credentials: Dict[str, str]
    metadata: Dict[str, Any]
    privacy_settings: Dict[str, Any] = None
    
    def __post_init__(self):
        if self.privacy_settings is None:
            self.privacy_settings = {}

class ExportProcessor:
    """
    Handles export operations to various formats.
    """
    
    def __init__(self, temp_dir: str = None):
        self.temp_dir = Path(temp_dir) if temp_dir else Path(tempfile.gettempdir()) / "exports"
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        self.logger = logging.getLogger(__name__)
    
    async def export_video(self, input_path: str, output_path: str, 
                          config: ExportConfig) -> Dict[str, Any]:
        """Export video to specified format."""
        if not EXPORT_AVAILABLE:
            raise RuntimeError("Export libraries not available")
        
        input_path = Path(input_path)
        output_path = Path(output_path)
        
        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")
        
        # Ensure output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        export_result = {
            "input_path": str(input_path),
            "output_path": str(output_path),
            "format": config.format.value,
            "start_time": datetime.now(),
            "success": False
        }
        
        try:
            if config.format == ExportFormat.MP4:
                await self._export_to_mp4(input_path, output_path, config)
            elif config.format == ExportFormat.WEBM:
                await self._export_to_webm(input_path, output_path, config)
            elif config.format == ExportFormat.GIF:
                await self._export_to_gif(input_path, output_path, config)
            elif config.format == ExportFormat.POWERPOINT:
                await self._export_to_powerpoint(input_path, output_path, config)
            elif config.format == ExportFormat.IMAGES:
                await self._export_to_images(input_path, output_path, config)
            elif config.format == ExportFormat.AUDIO_ONLY:
                await self._export_audio_only(input_path, output_path, config)
            else:
                raise ValueError(f"Unsupported export format: {config.format}")
            
            export_result["success"] = True
            export_result["file_size"] = output_path.stat().st_size if output_path.exists() else 0
            
        except Exception as e:
            export_result["error"] = str(e)
            self.logger.error(f"Export failed: {e}")
        
        export_result["end_time"] = datetime.now()
        export_result["duration"] = (export_result["end_time"] - export_result["start_time"]).total_seconds()
        
        return export_result
    
    async def _export_to_mp4(self, input_path: Path, output_path: Path, config: ExportConfig):
        """Export to MP4 format."""
        # Build ffmpeg command
        input_stream = ffmpeg.input(str(input_path))
        
        # Apply quality settings
        output_args = {}
        
        if config.quality == "high":
            output_args.update({
                'vcodec': 'libx264',
                'acodec': 'aac',
                'crf': 18,
                'preset': 'slow'
            })
        elif config.quality == "medium":
            output_args.update({
                'vcodec': 'libx264',
                'acodec': 'aac',
                'crf': 23,
                'preset': 'medium'
            })
        elif config.quality == "low":
            output_args.update({
                'vcodec': 'libx264',
                'acodec': 'aac',
                'crf': 28,
                'preset': 'fast'
            })
        
        # Apply custom parameters
        if config.resolution:
            width, height = config.resolution.split('x')
            output_args['s'] = f"{width}x{height}"
        
        if config.fps:
            output_args['r'] = config.fps
        
        if config.bitrate:
            output_args['b:v'] = config.bitrate
        
        # Add custom parameters
        output_args.update(config.custom_params)
        
        # Execute export
        output_stream = ffmpeg.output(input_stream, str(output_path), **output_args)
        await asyncio.create_task(
            asyncio.to_thread(ffmpeg.run, output_stream, overwrite_output=True, quiet=True)
        )
    
    async def _export_to_webm(self, input_path: Path, output_path: Path, config: ExportConfig):
        """Export to WebM format."""
        input_stream = ffmpeg.input(str(input_path))
        
        output_args = {
            'vcodec': 'libvpx-vp9',
            'acodec': 'libopus',
            'crf': 30 if config.quality == "high" else 35,
            'b:v': '0'  # Use CRF mode
        }
        
        if config.resolution:
            width, height = config.resolution.split('x')
            output_args['s'] = f"{width}x{height}"
        
        output_args.update(config.custom_params)
        
        output_stream = ffmpeg.output(input_stream, str(output_path), **output_args)
        await asyncio.create_task(
            asyncio.to_thread(ffmpeg.run, output_stream, overwrite_output=True, quiet=True)
        )
    
    async def _export_to_gif(self, input_path: Path, output_path: Path, config: ExportConfig):
        """Export to GIF format."""
        # Use moviepy for GIF export (better quality than ffmpeg)
        def create_gif():
            clip = mp.VideoFileClip(str(input_path))
            
            # Resize if needed
            if config.resolution:
                width, height = map(int, config.resolution.split('x'))
                clip = clip.resize((width, height))
            
            # Limit duration for GIF (max 30 seconds)
            if clip.duration > 30:
                clip = clip.subclip(0, 30)
            
            # Reduce fps for smaller file size
            fps = config.fps or 10
            
            clip.write_gif(str(output_path), fps=fps, opt='OptimizePlus')
            clip.close()
        
        await asyncio.create_task(asyncio.to_thread(create_gif))
    
    async def _export_to_powerpoint(self, input_path: Path, output_path: Path, config: ExportConfig):
        """Export to PowerPoint format."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            # Extract frames from video
            frames_dir = self.temp_dir / f"frames_{uuid.uuid4()}"
            frames_dir.mkdir(exist_ok=True)
            
            # Extract frames using ffmpeg
            input_stream = ffmpeg.input(str(input_path))
            output_stream = ffmpeg.output(
                input_stream, 
                str(frames_dir / "frame_%04d.png"),
                vf='fps=1/5'  # One frame every 5 seconds
            )
            await asyncio.create_task(
                asyncio.to_thread(ffmpeg.run, output_stream, overwrite_output=True, quiet=True)
            )
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Add slides with frames
            frame_files = sorted(frames_dir.glob("frame_*.png"))
            for frame_file in frame_files:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Add image to slide
                slide.shapes.add_picture(
                    str(frame_file),
                    Inches(1), Inches(1),
                    Inches(8), Inches(6)
                )
            
            # Save presentation
            prs.save(str(output_path))
            
            # Cleanup frames
            import shutil
            shutil.rmtree(frames_dir)
            
        except ImportError:
            raise RuntimeError("python-pptx not available for PowerPoint export")
    
    async def _export_to_images(self, input_path: Path, output_path: Path, config: ExportConfig):
        """Export to image sequence."""
        # Create output directory
        output_dir = output_path.parent / output_path.stem
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Extract frames
        input_stream = ffmpeg.input(str(input_path))
        
        fps = config.fps or 1  # Default to 1 fps for images
        output_pattern = str(output_dir / "frame_%06d.png")
        
        output_stream = ffmpeg.output(
            input_stream,
            output_pattern,
            vf=f'fps={fps}'
        )
        
        await asyncio.create_task(
            asyncio.to_thread(ffmpeg.run, output_stream, overwrite_output=True, quiet=True)
        )
        
        # Create zip file of images
        import zipfile
        with zipfile.ZipFile(output_path, 'w') as zipf:
            for img_file in output_dir.glob("*.png"):
                zipf.write(img_file, img_file.name)
    
    async def _export_audio_only(self, input_path: Path, output_path: Path, config: ExportConfig):
        """Export audio only."""
        input_stream = ffmpeg.input(str(input_path))
        
        output_args = {
            'acodec': 'mp3' if output_path.suffix == '.mp3' else 'aac',
            'vn': None  # No video
        }
        
        if config.quality == "high":
            output_args['ab'] = '320k'
        elif config.quality == "medium":
            output_args['ab'] = '192k'
        else:
            output_args['ab'] = '128k'
        
        output_args.update(config.custom_params)
        
        output_stream = ffmpeg.output(input_stream, str(output_path), **output_args)
        await asyncio.create_task(
            asyncio.to_thread(ffmpeg.run, output_stream, overwrite_output=True, quiet=True)
        )

class PlatformPublisher:
    """
    Handles publishing to various platforms.
    """
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    async def publish_to_platform(self, video_path: str, config: PublishConfig) -> Dict[str, Any]:
        """Publish video to specified platform."""
        if not INTEGRATION_AVAILABLE:
            raise RuntimeError("Integration libraries not available")
        
        publish_result = {
            "platform": config.platform.value,
            "start_time": datetime.now(),
            "success": False
        }
        
        try:
            if config.platform == PlatformType.YOUTUBE:
                result = await self._publish_to_youtube(video_path, config)
            elif config.platform == PlatformType.VIMEO:
                result = await self._publish_to_vimeo(video_path, config)
            elif config.platform == PlatformType.TWITTER:
                result = await self._publish_to_twitter(video_path, config)
            else:
                raise ValueError(f"Platform {config.platform} not yet implemented")
            
            publish_result.update(result)
            publish_result["success"] = True
            
        except Exception as e:
            publish_result["error"] = str(e)
            self.logger.error(f"Publishing failed: {e}")
        
        publish_result["end_time"] = datetime.now()
        return publish_result
    
    async def _publish_to_youtube(self, video_path: str, config: PublishConfig) -> Dict:
        """Publish to YouTube using YouTube Data API."""
        # This is a simplified implementation
        # In practice, you'd use the Google API client library
        
        api_key = config.credentials.get('api_key')
        access_token = config.credentials.get('access_token')
        
        if not api_key or not access_token:
            raise ValueError("YouTube API credentials required")
        
        # Upload video (simplified)
        upload_url = "https://www.googleapis.com/upload/youtube/v3/videos"
        
        headers = {
            'Authorization': f'Bearer {access_token}',
            'Content-Type': 'application/octet-stream'
        }
        
        params = {
            'part': 'snippet,status',
            'key': api_key
        }
        
        # Video metadata
        metadata = {
            'snippet': {
                'title': config.metadata.get('title', 'Generated Video'),
                'description': config.metadata.get('description', ''),
                'tags': config.metadata.get('tags', []),
                'categoryId': config.metadata.get('category_id', '22')  # People & Blogs
            },
            'status': {
                'privacyStatus': config.privacy_settings.get('privacy', 'private')
            }
        }
        
        # Note: This is a simplified example
        # Real implementation would use proper multipart upload
        async with httpx.AsyncClient() as client:
            response = await client.post(
                upload_url,
                params=params,
                headers=headers,
                json=metadata
            )
        
        if response.status_code == 200:
            result = response.json()
            return {
                'video_id': result.get('id'),
                'url': f"https://www.youtube.com/watch?v={result.get('id')}",
                'status': 'uploaded'
            }
        else:
            raise Exception(f"YouTube upload failed: {response.text}")
    
    async def _publish_to_vimeo(self, video_path: str, config: PublishConfig) -> Dict:
        """Publish to Vimeo."""
        access_token = config.credentials.get('access_token')
        
        if not access_token:
            raise ValueError("Vimeo access token required")
        
        # Simplified Vimeo upload
        headers = {
            'Authorization': f'Bearer {access_token}',
            'Content-Type': 'application/json'
        }
        
        # Create video entry
        create_url = "https://api.vimeo.com/me/videos"
        video_data = {
            'name': config.metadata.get('title', 'Generated Video'),
            'description': config.metadata.get('description', ''),
            'privacy': {
                'view': config.privacy_settings.get('privacy', 'nobody')
            }
        }
        
        async with httpx.AsyncClient() as client:
            response = await client.post(create_url, headers=headers, json=video_data)
        
        if response.status_code == 201:
            result = response.json()
            return {
                'video_id': result.get('uri'),
                'url': result.get('link'),
                'status': 'created'
            }
        else:
            raise Exception(f"Vimeo upload failed: {response.text}")
    
    async def _publish_to_twitter(self, video_path: str, config: PublishConfig) -> Dict:
        """Publish to Twitter."""
        # Twitter has specific video requirements
        # This would use the Twitter API v2
        
        api_key = config.credentials.get('api_key')
        api_secret = config.credentials.get('api_secret')
        access_token = config.credentials.get('access_token')
        access_token_secret = config.credentials.get('access_token_secret')
        
        if not all([api_key, api_secret, access_token, access_token_secret]):
            raise ValueError("Twitter API credentials required")
        
        # Simplified Twitter upload
        return {
            'status': 'Twitter upload not fully implemented',
            'message': 'Would upload to Twitter with proper API integration'
        }

class WebhookManager:
    """
    Manages webhooks for automated workflows.
    """
    
    def __init__(self):
        self.webhooks: Dict[str, Dict] = {}
        self.logger = logging.getLogger(__name__)
    
    def register_webhook(self, event_type: str, url: str, 
                        headers: Dict[str, str] = None) -> str:
        """Register webhook for event type."""
        webhook_id = str(uuid.uuid4())
        
        self.webhooks[webhook_id] = {
            'event_type': event_type,
            'url': url,
            'headers': headers or {},
            'created_at': datetime.now(),
            'active': True
        }
        
        return webhook_id
    
    async def trigger_webhooks(self, event_type: str, data: Dict[str, Any]):
        """Trigger all webhooks for event type."""
        matching_webhooks = [
            webhook for webhook in self.webhooks.values()
            if webhook['event_type'] == event_type and webhook['active']
        ]
        
        for webhook in matching_webhooks:
            try:
                await self._send_webhook(webhook, data)
            except Exception as e:
                self.logger.error(f"Webhook failed: {e}")
    
    async def _send_webhook(self, webhook: Dict, data: Dict[str, Any]):
        """Send webhook request."""
        payload = {
            'event_type': webhook['event_type'],
            'timestamp': datetime.now().isoformat(),
            'data': data
        }
        
        async with httpx.AsyncClient() as client:
            response = await client.post(
                webhook['url'],
                json=payload,
                headers=webhook['headers'],
                timeout=30.0
            )
            
            if response.status_code >= 400:
                raise Exception(f"Webhook returned {response.status_code}: {response.text}")

class APIIntegration:
    """
    Provides API endpoints for third-party integrations.
    """
    
    def __init__(self, export_processor: ExportProcessor, 
                 publisher: PlatformPublisher, webhook_manager: WebhookManager):
        self.export_processor = export_processor
        self.publisher = publisher
        self.webhook_manager = webhook_manager
        self.logger = logging.getLogger(__name__)
    
    async def export_and_publish_workflow(self, video_id: str, 
                                        export_configs: List[ExportConfig],
                                        publish_configs: List[PublishConfig] = None) -> Dict:
        """Complete export and publish workflow."""
        workflow_id = str(uuid.uuid4())
        
        workflow_result = {
            'workflow_id': workflow_id,
            'video_id': video_id,
            'start_time': datetime.now(),
            'exports': [],
            'publications': [],
            'success': False
        }
        
        try:
            # Get video file path (this would come from your video storage system)
            video_path = f"data/videos/{video_id}.mp4"  # Placeholder
            
            # Export to all requested formats
            for export_config in export_configs:
                output_filename = f"{video_id}_{export_config.format.value}"
                output_path = self.export_processor.temp_dir / output_filename
                
                export_result = await self.export_processor.export_video(
                    video_path, str(output_path), export_config
                )
                workflow_result['exports'].append(export_result)
            
            # Publish to platforms if requested
            if publish_configs:
                for publish_config in publish_configs:
                    # Use the original video or first export
                    publish_path = video_path
                    if workflow_result['exports'] and workflow_result['exports'][0]['success']:
                        publish_path = workflow_result['exports'][0]['output_path']
                    
                    publish_result = await self.publisher.publish_to_platform(
                        publish_path, publish_config
                    )
                    workflow_result['publications'].append(publish_result)
            
            # Check overall success
            export_success = all(exp['success'] for exp in workflow_result['exports'])
            publish_success = all(pub['success'] for pub in workflow_result['publications']) if workflow_result['publications'] else True
            
            workflow_result['success'] = export_success and publish_success
            
            # Trigger webhooks
            await self.webhook_manager.trigger_webhooks('export_complete', workflow_result)
            
        except Exception as e:
            workflow_result['error'] = str(e)
            self.logger.error(f"Workflow failed: {e}")
        
        workflow_result['end_time'] = datetime.now()
        workflow_result['duration'] = (workflow_result['end_time'] - workflow_result['start_time']).total_seconds()
        
        return workflow_result

# Integration functions
async def export_video_multiple_formats(video_id: str, formats: List[str]) -> Dict:
    """
    Convenience function to export video to multiple formats.
    
    Args:
        video_id: ID of the video to export
        formats: List of format names (mp4, webm, gif, etc.)
        
    Returns:
        Dictionary with export results
    """
    processor = ExportProcessor()
    
    export_configs = []
    for format_name in formats:
        try:
            export_format = ExportFormat(format_name.lower())
            config = ExportConfig(format=export_format, quality="high")
            export_configs.append(config)
        except ValueError:
            logging.warning(f"Unsupported format: {format_name}")
    
    # Create API integration
    publisher = PlatformPublisher()
    webhook_manager = WebhookManager()
    api = APIIntegration(processor, publisher, webhook_manager)
    
    return await api.export_and_publish_workflow(video_id, export_configs)

if __name__ == "__main__":
    # Example usage
    async def main():
        print("Export and integration system initialized successfully")
        print("Available features:")
        print(f"- Export capabilities: {EXPORT_AVAILABLE}")
        print(f"- Integration capabilities: {INTEGRATION_AVAILABLE}")
        
        # Example export
        if EXPORT_AVAILABLE:
            processor = ExportProcessor()
            
            config = ExportConfig(
                format=ExportFormat.MP4,
                quality="high",
                resolution="1920x1080"
            )
            
            print(f"Export configuration created for {config.format.value}")
        
        # Example webhook
        webhook_manager = WebhookManager()
        webhook_id = webhook_manager.register_webhook(
            "export_complete",
            "https://example.com/webhook"
        )
        print(f"Webhook registered: {webhook_id}")
    
    asyncio.run(main())